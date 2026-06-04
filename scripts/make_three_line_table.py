#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Create editable three-line tables in a PPTX using python-pptx.

This helper can add a new slide containing a CUMT-style three-line table, or it
can be imported and the `add_three_line_table` function reused by deck builders.
"""

from __future__ import annotations

import argparse
import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BLUE = RGBColor(0x1E, 0x3C, 0x6B)
DARK = RGBColor(0x1F, 0x29, 0x37)
LIGHT = RGBColor(0xF7, 0xFA, 0xFD)
FONT_CN_TITLE = "SimHei"
FONT_CN_BODY = "SimSun"
FONT_EN = "Times New Roman"


def is_cjk(ch: str) -> bool:
    return "\u4e00" <= ch <= "\u9fff" or "\u3000" <= ch <= "\u303f" or "\uff00" <= ch <= "\uffef"


def ensure_child(parent, tag):
    child = parent.find(qn(tag))
    if child is None:
        child = OxmlElement(tag)
        parent.append(child)
    return child


def set_typeface(run, font):
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        child = ensure_child(rpr, tag)
        child.set("typeface", font)


def add_mixed_text(slide, x, y, w, h, text, size=15, bold=False, header=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    current = ""
    current_font = None
    for ch in text:
        font = (FONT_CN_TITLE if header else FONT_CN_BODY) if is_cjk(ch) else FONT_EN
        if current and font != current_font:
            run = p.add_run()
            run.text = current
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = BLUE if header else DARK
            set_typeface(run, current_font)
            current = ch
            current_font = font
        else:
            current += ch
            current_font = font
    if current:
        run = p.add_run()
        run.text = current
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = BLUE if header else DARK
        set_typeface(run, current_font)
    return box


def add_rule(slide, x, y, w, color=BLUE, points=1.5):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.height = Pt(points)
    return line


def add_three_line_table(slide, rows, x=1.2, y=1.4, widths=(2.0, 2.5, 6.0), row_h=0.45):
    """Add an editable three-line table.

    `rows` should include the header row as the first row.
    """
    total_w = sum(widths)
    add_rule(slide, x, y, total_w, BLUE, 1.8)
    add_rule(slide, x, y + row_h, total_w, BLUE, 1.4)
    add_rule(slide, x, y + row_h * len(rows), total_w, BLUE, 1.8)
    for i, row in enumerate(rows):
        xx = x
        for j, cell in enumerate(row):
            add_mixed_text(
                slide,
                xx + 0.08,
                y + i * row_h + 0.08,
                widths[j] - 0.16,
                row_h - 0.12,
                str(cell),
                size=15 if i else 16,
                bold=i == 0,
                header=i == 0,
                align=PP_ALIGN.CENTER if j < 2 else PP_ALIGN.LEFT,
            )
            xx += widths[j]


def read_csv(path: Path):
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.reader(f))


def main() -> int:
    parser = argparse.ArgumentParser(description="Add a three-line table slide to a PPTX.")
    parser.add_argument("input_pptx", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--csv", type=Path, required=True, help="CSV rows, first row is header.")
    args = parser.parse_args()

    prs = Presentation(str(args.input_pptx))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rows = read_csv(args.csv)
    add_three_line_table(slide, rows)
    prs.save(str(args.output_pptx))
    print(f"output={args.output_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
