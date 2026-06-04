#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Normalize upper-right CUMT logo size and position.

The script uses a standard slide's upper-right logo as the source of truth and
copies its left/top/width/height to logos on selected later slides.
"""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def logo_candidates(prs, slide, image_size=None):
    candidates = []
    for shape in iter_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.left < prs.slide_width * 0.65 or shape.top > prs.slide_height * 0.25:
            continue
        if image_size is not None and shape.image.size != tuple(image_size):
            continue
        candidates.append(shape)
    return candidates


def main() -> int:
    parser = argparse.ArgumentParser(description="Normalize CUMT logo placement.")
    parser.add_argument("input_pptx", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--standard-slide", type=int, default=3, help="1-based slide number used as logo standard.")
    parser.add_argument("--start-slide", type=int, default=4, help="1-based first slide to normalize.")
    parser.add_argument("--image-size", default=None, help="Optional natural image size filter, e.g. 424x108.")
    args = parser.parse_args()

    image_size = None
    if args.image_size:
        w, h = args.image_size.lower().split("x")
        image_size = (int(w), int(h))

    shutil.copy2(args.input_pptx, args.output_pptx)
    prs = Presentation(str(args.output_pptx))
    standard_slide = prs.slides[args.standard_slide - 1]
    standards = logo_candidates(prs, standard_slide, image_size)
    if len(standards) != 1:
        raise RuntimeError(f"Expected one logo on standard slide {args.standard_slide}, found {len(standards)}.")
    std = standards[0]

    changed = []
    unidentified = []
    for slide_no in range(args.start_slide, len(prs.slides) + 1):
        candidates = logo_candidates(prs, prs.slides[slide_no - 1], image_size or std.image.size)
        if len(candidates) != 1:
            unidentified.append(slide_no)
            continue
        logo = candidates[0]
        logo.left = std.left
        logo.top = std.top
        logo.width = std.width
        logo.height = std.height
        changed.append(slide_no)

    prs.save(str(args.output_pptx))
    print(f"output={args.output_pptx}")
    print(
        "standard="
        f"slide {args.standard_slide}, left={std.left}, top={std.top}, "
        f"width={std.width}, height={std.height}"
    )
    print(f"changed={changed}")
    print(f"unidentified={unidentified}")
    return 0 if not unidentified else 2


if __name__ == "__main__":
    raise SystemExit(main())
