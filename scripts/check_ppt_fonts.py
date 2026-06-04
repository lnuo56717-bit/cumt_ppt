#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Check common CUMT thesis-defense PPT font rules.

This is a heuristic checker for python-pptx. It reports likely violations but
does not modify the file.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


TITLE_CN = "SimHei"
BODY_CN = "SimSun"
EN = "Times New Roman"


def is_cjk(ch: str) -> bool:
    code = ord(ch)
    return (
        0x4E00 <= code <= 0x9FFF
        or 0x3400 <= code <= 0x4DBF
        or 0xF900 <= code <= 0xFAFF
        or 0x3000 <= code <= 0x303F
        or 0xFF00 <= code <= 0xFFEF
    )


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(
        "".join(run.text for run in paragraph.runs)
        for paragraph in shape.text_frame.paragraphs
    )


def role_for(shape, slide_index: int) -> str:
    text = shape_text(shape).strip()
    first = text.splitlines()[0].strip() if text else ""
    if not first:
        return "empty"
    if slide_index == 0 and len(first) >= 8:
        return "title"
    if first in {"致谢", "请批评指正", "毕业设计答辩"}:
        return "title"
    if re.match(r"^\d{2}\s+", first) and shape.top < 1_000_000:
        return "title"
    first_size = None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size:
                first_size = run.font.size.pt
                break
        if first_size:
            break
    if first_size and first_size >= 20 and len(first) <= 18:
        return "module"
    if first_size and first_size <= 14:
        return "small"
    return "body"


def expected_font(text: str, role: str) -> str:
    if any(is_cjk(ch) for ch in text):
        return TITLE_CN if role in {"title", "module"} else BODY_CN
    return EN


def check(pptx: Path) -> dict:
    prs = Presentation(str(pptx))
    issues = []
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape_index, shape in enumerate(iter_shapes(slide.shapes), 1):
            if not getattr(shape, "has_text_frame", False):
                continue
            role = role_for(shape, slide_index - 1)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if not text.strip():
                        continue
                    expected = expected_font(text, role)
                    actual = run.font.name
                    if actual != expected:
                        issues.append(
                            {
                                "slide": slide_index,
                                "shape": shape_index,
                                "role": role,
                                "text": text,
                                "actual": actual,
                                "expected": expected,
                            }
                        )
    return {
        "pptx": str(pptx),
        "slides": len(prs.slides),
        "issue_count": len(issues),
        "issues": issues,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Check CUMT PPT font rules.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--json", action="store_true", help="Print full JSON report.")
    args = parser.parse_args()

    report = check(args.pptx)
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"PPTX: {report['pptx']}")
        print(f"Slides: {report['slides']}")
        print(f"Font issues: {report['issue_count']}")
        for issue in report["issues"][:50]:
            print(
                f"slide {issue['slide']} shape {issue['shape']}: "
                f"{issue['text']!r} actual={issue['actual']!r} expected={issue['expected']!r}"
            )
        if report["issue_count"] > 50:
            print(f"... {report['issue_count'] - 50} more")
    return 0 if report["issue_count"] == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
