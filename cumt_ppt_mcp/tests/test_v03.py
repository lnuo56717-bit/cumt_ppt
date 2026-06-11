from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
sys.path.insert(0, str(SRC))

from cumt_ppt_mcp.ppt_utils import (  # noqa: E402
    add_cumt_logo,
    apply_template_style,
    auto_fix_layout,
    check_cumt_logo_asset,
    check_layout_overlap,
    extract_pdf_images_safe,
    extract_logo_style,
    inspect_template_style,
    is_logo_candidate,
    normalize_logo,
)


LOGO = ROOT / "assets" / "cumt_logo.png"


def _sample_deck(path: Path, overlap: bool = False) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for i in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(7.5), Inches(0.5))
        run = title.text_frame.paragraphs[0].add_run()
        run.text = f"0{i + 1} 模板测试标题"
        run.font.size = Pt(28)
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(5.0), Inches(1.5))
        body.text = "这是虚构测试正文，用于检查布局和字体。"
        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), Inches(11.0), Inches(0.1), width=Inches(1.7))
            if overlap and i == 1:
                slide.shapes.add_picture(str(LOGO), Inches(1.0), Inches(1.2), width=Inches(4.8))
    prs.save(path)
    return path


def _template_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.62))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(29, 67, 113)
        bar.line.color.rgb = RGBColor(29, 67, 113)
        title = slide.shapes.add_textbox(Inches(0.9), Inches(0.12), Inches(8.0), Inches(0.35))
        run = title.text_frame.paragraphs[0].add_run()
        run.text = f"0{i + 1} 模板标题"
        run.font.size = Pt(22)
        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), Inches(10.9), Inches(0.1), width=Inches(1.7))
    prs.save(path)
    return path


def _complex_template_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    content_image = Image.new("RGB", (900, 420), (235, 240, 248))
    for x in range(0, 900, 60):
        for y in range(0, 420, 60):
            content_image.putpixel((x, y), (40, 90, 150))
    content_path = path.with_name("content_image.png")
    content_image.save(content_path)
    deco_image = Image.new("RGB", (320, 40), (210, 30, 30))
    deco_path = path.with_name("bottom_deco.png")
    deco_image.save(deco_path)

    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.62))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(29, 67, 113)
        bar.line.color.rgb = RGBColor(29, 67, 113)
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.14), Inches(6.5), Inches(0.36))
        run = title.text_frame.paragraphs[0].add_run()
        run.text = f"0{i + 1} complex template"
        run.font.size = Pt(22)
        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), Inches(11.15), Inches(0.1), width=Inches(1.45))
        slide.shapes.add_picture(str(content_path), Inches(2.0), Inches(1.4), width=Inches(6.8))
        slide.shapes.add_picture(str(deco_path), Inches(5.0), Inches(6.85), width=Inches(2.6))
    prs.save(path)
    return path


def test_logo_asset_and_logo_tools(tmp_path: Path) -> None:
    asset = check_cumt_logo_asset()
    assert asset["ok"] is True
    assert asset["size_px"]["width"] >= 200
    assert asset["black_border_ratio"] < 0.65
    deck = _sample_deck(tmp_path / "deck.pptx")
    added = tmp_path / "added.pptx"
    result = add_cumt_logo(str(deck), str(added), slide_indices=[2, 3], skip_slide_indices=[2])
    assert result["ok"] is True
    assert added.exists()
    normalized = tmp_path / "normalized.pptx"
    norm = normalize_logo(str(added), str(normalized), reference_slide_index=3, start_slide_index=4)
    assert norm["ok"] is True
    assert normalized.exists()


def test_template_style_and_apply(tmp_path: Path) -> None:
    source = _sample_deck(tmp_path / "source.pptx")
    template = _template_deck(tmp_path / "template.pptx")
    style = inspect_template_style(str(template))
    assert style["ok"] is True
    assert style["title_bar_candidates"]
    out = tmp_path / "styled.pptx"
    applied = apply_template_style(str(source), str(template), str(out))
    assert applied["ok"] is True
    assert out.exists()


def test_template_image_classification_avoids_content_logo(tmp_path: Path) -> None:
    source = _sample_deck(tmp_path / "source.pptx")
    template = _complex_template_deck(tmp_path / "complex_template.pptx")
    prs = Presentation(str(template))
    first_slide_pictures = [shape for shape in prs.slides[0].shapes if shape.shape_type == 13]
    assert len(first_slide_pictures) == 3
    assert is_logo_candidate(first_slide_pictures[0], prs.slide_width, prs.slide_height)
    assert not is_logo_candidate(first_slide_pictures[1], prs.slide_width, prs.slide_height)
    assert not is_logo_candidate(first_slide_pictures[2], prs.slide_width, prs.slide_height)

    style = inspect_template_style(str(template))
    image_classes = [item["classification"] for item in style["image_classifications"]]
    assert image_classes.count("header_logo") == 3
    assert image_classes.count("content_image") >= 3
    assert any(item["classification"] == "background_decoration" for item in style["image_classifications"])

    logo_style = extract_logo_style(str(template))
    assert len(logo_style["header_logos"]) == 3

    out = tmp_path / "styled_no_big_logo.pptx"
    result = apply_template_style(str(source), str(template), str(out))
    assert result["ok"] is True
    assert out.exists()
    styled = Presentation(str(out))
    for idx, slide in enumerate(styled.slides):
        if idx == 0:
            continue
        pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
        assert all(shape.width < styled.slide_width * 0.20 for shape in pictures)


def test_layout_overlap_and_auto_fix(tmp_path: Path) -> None:
    deck = _sample_deck(tmp_path / "overlap.pptx", overlap=True)
    report = tmp_path / "layout.md"
    checked = check_layout_overlap(str(deck), str(report))
    assert checked["ok"] is True
    assert checked["issue_count"] >= 1
    fixed = tmp_path / "fixed.pptx"
    fix = auto_fix_layout(str(deck), str(fixed))
    assert fix["ok"] is True
    assert fixed.exists()


def test_extract_pdf_images_safe(tmp_path: Path) -> None:
    import fitz

    pdf = tmp_path / "transparent_image.pdf"
    png = tmp_path / "alpha.png"
    img = Image.new("RGBA", (240, 120), (0, 0, 0, 0))
    for x in range(40, 200):
        for y in range(30, 90):
            img.putpixel((x, y), (30, 90, 160, 180))
    img.save(png)
    doc = fitz.open()
    page = doc.new_page(width=400, height=260)
    page.insert_image(fitz.Rect(60, 60, 300, 180), filename=str(png))
    doc.save(pdf)
    out_dir = tmp_path / "pdf_out"
    result = extract_pdf_images_safe(str(pdf), str(out_dir))
    assert result["ok"] is True
    assert result["rendered_page_count"] == 1
    assert Path(result["report_path"]).exists()
