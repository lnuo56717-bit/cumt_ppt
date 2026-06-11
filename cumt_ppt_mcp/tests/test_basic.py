from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
sys.path.insert(0, str(SRC))

from cumt_ppt_mcp.ppt_utils import (  # noqa: E402
    apply_font_rules,
    check_ppt_quality,
    inspect_ppt,
    inspect_slide,
    make_three_line_table,
    normalize_logo,
)


def _make_sample_pptx(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for idx in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.5), Inches(0.5))
        p = title.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"0{idx + 1} 测试标题 AGV {idx + 1}"
        run.font.size = Pt(28)
        body = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(7.5), Inches(1.0))
        body.text = "正文内容 PPO 147 22 32428.0"
        # Use a simple rectangle saved as an in-deck picture-like logo is hard without image files;
        # tests cover logo normalization only when a fixture picture exists.
    prs.save(path)
    return path


def test_inspect_ppt_and_slide(tmp_path: Path) -> None:
    pptx = _make_sample_pptx(tmp_path / "sample.pptx")
    info = inspect_ppt(str(pptx))
    assert info["slide_count"] == 4
    assert info["slides"][0]["text_box_count"] >= 2
    slide = inspect_slide(str(pptx), 1)
    assert slide["slide_index"] == 1
    assert "测试标题" in slide["title"]


def test_apply_font_rules_outputs_new_file(tmp_path: Path) -> None:
    pptx = _make_sample_pptx(tmp_path / "sample.pptx")
    out = tmp_path / "fonts.pptx"
    result = apply_font_rules(str(pptx), str(out))
    assert result["ok"] is True
    assert out.exists()
    assert out != pptx


def test_make_three_line_table_and_quality_report(tmp_path: Path) -> None:
    pptx = _make_sample_pptx(tmp_path / "sample.pptx")
    out = tmp_path / "table.pptx"
    result = make_three_line_table(
        str(pptx),
        str(out),
        2,
        [["参数类别", "参数名称", "参数取值"], ["算法", "优化算法", "SMGOA"], ["结果", "总延迟", "1210"]],
    )
    assert result["ok"] is True
    assert out.exists()
    quality = check_ppt_quality(str(out), expected_title="测试标题", required_numbers=["1210"])
    assert quality["ok"] is True
    assert Path(quality["report_path"]).exists()


def test_reject_overwriting_input(tmp_path: Path) -> None:
    pptx = _make_sample_pptx(tmp_path / "sample.pptx")
    result_path = str(pptx)
    with pytest.raises(Exception):
        apply_font_rules(str(pptx), result_path)
