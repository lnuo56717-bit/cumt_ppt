from __future__ import annotations

from typing import Any

from mcp.server.fastmcp import FastMCP

from .ppt_utils import (
    PptToolError,
    add_cumt_logo as add_cumt_logo_impl,
    apply_font_rules as apply_font_rules_impl,
    apply_template_style as apply_template_style_impl,
    auto_fix_layout as auto_fix_layout_impl,
    check_cumt_logo_asset as check_cumt_logo_asset_impl,
    check_layout_overlap as check_layout_overlap_impl,
    check_ppt_quality as check_ppt_quality_impl,
    export_preview as export_preview_impl,
    extract_pdf_images_safe as extract_pdf_images_safe_impl,
    extract_logo_style as extract_logo_style_impl,
    inspect_ppt as inspect_ppt_impl,
    inspect_slide as inspect_slide_impl,
    inspect_template_style as inspect_template_style_impl,
    make_three_line_table as make_three_line_table_impl,
    normalize_logo as normalize_logo_impl,
    resize_images_keep_ratio as resize_images_keep_ratio_impl,
)


mcp = FastMCP("cumt_ppt_mcp")


def _safe_call(fn, *args, **kwargs) -> dict[str, Any]:
    try:
        return fn(*args, **kwargs)
    except PptToolError as exc:
        return {"ok": False, "error": str(exc)}
    except Exception as exc:
        return {"ok": False, "error": f"{type(exc).__name__}: {exc}"}


@mcp.tool()
def inspect_ppt(pptx_path: str) -> dict[str, Any]:
    """Inspect a PPTX deck: slide count, titles, shape counts, and page size."""
    return _safe_call(inspect_ppt_impl, pptx_path)


@mcp.tool()
def inspect_slide(pptx_path: str, slide_index: int) -> dict[str, Any]:
    """Inspect one 1-based slide and return text, picture, table, and density details."""
    return _safe_call(inspect_slide_impl, pptx_path, slide_index)


@mcp.tool()
def apply_font_rules(
    pptx_path: str,
    output_path: str,
    title_font: str = "黑体",
    body_font: str = "宋体",
    latin_font: str = "Times New Roman",
) -> dict[str, Any]:
    """Apply CUMT defense font rules and save a new PPTX without changing text."""
    return _safe_call(apply_font_rules_impl, pptx_path, output_path, title_font, body_font, latin_font)


@mcp.tool()
def normalize_logo(
    pptx_path: str,
    output_path: str,
    reference_slide_index: int = 3,
    start_slide_index: int = 4,
    skip_slide_indices: list[int] | None = None,
    logo_asset_path: str | None = None,
) -> dict[str, Any]:
    """Normalize right-top CUMT logos using the fixed local logo asset."""
    return _safe_call(
        normalize_logo_impl,
        pptx_path,
        output_path,
        reference_slide_index,
        start_slide_index,
        skip_slide_indices,
        logo_asset_path,
    )


@mcp.tool()
def check_cumt_logo_asset(logo_asset_path: str | None = None) -> dict[str, Any]:
    """Check the reusable CUMT logo asset for black background, alpha, and size risks."""
    return _safe_call(check_cumt_logo_asset_impl, logo_asset_path)


@mcp.tool()
def add_cumt_logo(
    pptx_path: str,
    output_path: str,
    slide_indices: list[int] | None = None,
    logo_asset_path: str | None = None,
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    skip_slide_indices: list[int] | None = None,
) -> dict[str, Any]:
    """Add the fixed local CUMT logo asset to selected slides."""
    return _safe_call(add_cumt_logo_impl, pptx_path, output_path, slide_indices, logo_asset_path, left, top, width, skip_slide_indices)


@mcp.tool()
def resize_images_keep_ratio(
    pptx_path: str,
    output_path: str,
    slide_indices: list[int] | None = None,
    mode: str = "fit_content_area",
    max_width_ratio: float = 0.45,
    max_height_ratio: float = 0.58,
) -> dict[str, Any]:
    """Resize selected slide pictures while preserving aspect ratio."""
    return _safe_call(
        resize_images_keep_ratio_impl,
        pptx_path,
        output_path,
        slide_indices,
        mode,
        max_width_ratio,
        max_height_ratio,
    )


@mcp.tool()
def make_three_line_table(
    pptx_path: str,
    output_path: str,
    slide_index: int,
    table_data: list[list[Any]],
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    height: float | None = None,
) -> dict[str, Any]:
    """Add a thesis-style three-line table to one slide and save a new PPTX."""
    return _safe_call(make_three_line_table_impl, pptx_path, output_path, slide_index, table_data, left, top, width, height)


@mcp.tool()
def export_preview(pptx_path: str, output_dir: str) -> dict[str, Any]:
    """Export each slide to PNG using PowerPoint COM first, then LibreOffice."""
    return _safe_call(export_preview_impl, pptx_path, output_dir)


@mcp.tool()
def check_ppt_quality(
    pptx_path: str,
    expected_title: str | None = None,
    required_numbers: list[str] | None = None,
    report_path: str | None = None,
) -> dict[str, Any]:
    """Check title/numbers/placeholders/fonts/images/logo consistency and write a markdown report."""
    return _safe_call(check_ppt_quality_impl, pptx_path, expected_title, required_numbers, report_path)


@mcp.tool()
def inspect_template_style(template_pptx_path: str) -> dict[str, Any]:
    """Inspect a reference template PPTX and summarize style/layout features."""
    return _safe_call(inspect_template_style_impl, template_pptx_path)


@mcp.tool()
def extract_logo_style(template_pptx_path: str) -> dict[str, Any]:
    """Extract reusable header-logo placement only; content images are excluded."""
    return _safe_call(extract_logo_style_impl, template_pptx_path)


@mcp.tool()
def apply_template_style(source_pptx_path: str, template_pptx_path: str, output_path: str) -> dict[str, Any]:
    """Conservatively apply template title-bar, font, margin, and logo style without changing text."""
    return _safe_call(apply_template_style_impl, source_pptx_path, template_pptx_path, output_path)


@mcp.tool()
def check_layout_overlap(pptx_path: str, report_path: str | None = None) -> dict[str, Any]:
    """Check object overlap, out-of-bounds objects, and logo/title obstruction."""
    return _safe_call(check_layout_overlap_impl, pptx_path, report_path)


@mcp.tool()
def auto_fix_layout(pptx_path: str, output_path: str, report_path: str | None = None) -> dict[str, Any]:
    """Conservatively fix obvious image overlap/out-of-bounds issues."""
    return _safe_call(auto_fix_layout_impl, pptx_path, output_path, report_path)


@mcp.tool()
def extract_pdf_images_safe(pdf_path: str, output_dir: str, background_color: str = "white") -> dict[str, Any]:
    """Extract embedded PDF images and page renders with white-background alpha handling."""
    return _safe_call(extract_pdf_images_safe_impl, pdf_path, output_dir, background_color)


@mcp.tool()
def fit_image_in_box(img_path: str, box_w_inch: float, box_h_inch: float) -> dict[str, Any]:
    """Return display (width, height) in inches that fit an image inside a box preserving aspect ratio.

    Always call this before slide.shapes.add_picture() to prevent portrait images from
    overflowing the slide bottom. Returns: ok, width_inch, height_inch, scale, error.
    """
    from cumt_ppt_mcp.ppt_utils import fit_image_in_box as _impl
    return _impl(img_path, box_w_inch, box_h_inch)


@mcp.tool()
def generate_diagram(
    diagram_type: str,
    spec: dict[str, Any],
    output_path: str,
    figsize: list[float] | None = None,
    dpi: int = 160,
    title: str = "",
    lang: str = "zh",
) -> dict[str, Any]:
    """Generate a flowchart, layer diagram, or block diagram as a PNG using matplotlib.

    Covers diagram types that scipilot-figure-skill explicitly excludes (flowcharts,
    architecture diagrams, system block diagrams).

    Args:
        diagram_type: "flowchart" | "layer_diagram" | "block_diagram"
        spec: type-specific spec dict:
            flowchart    : {"nodes": [...], "edges": [...]}
            layer_diagram: {"layers": [...]}
            block_diagram: {"blocks": [...], "connections": [...]}
        output_path: output PNG path.
        figsize: [width, height] in inches (optional, auto-sized per type).
        dpi: output resolution (default 160).
        title: optional chart title.
        lang: "zh" or "en".

    Returns: dict with ok, output_path, error.
    """
    try:
        import sys
        import importlib.util
        from pathlib import Path

        script = Path(__file__).parent.parent.parent.parent.parent / "scripts" / "generate_diagram.py"
        spec_mod = importlib.util.spec_from_file_location("generate_diagram", str(script))
        if spec_mod is None:
            raise ImportError(f"Cannot find generate_diagram.py at {script}")
        mod = importlib.util.module_from_spec(spec_mod)
        spec_mod.loader.exec_module(mod)

        fs = tuple(figsize) if figsize else None
        kwargs = dict(output_path=output_path, dpi=dpi, title=title, lang=lang)
        if fs:
            kwargs["figsize"] = fs

        if diagram_type == "flowchart":
            mod.draw_flowchart(nodes=spec["nodes"], edges=spec["edges"], **kwargs)
        elif diagram_type == "layer_diagram":
            mod.draw_layer_diagram(layers=spec["layers"], **kwargs)
        elif diagram_type == "block_diagram":
            mod.draw_block_diagram(
                blocks=spec["blocks"], connections=spec["connections"], **kwargs)
        else:
            raise ValueError(f"Unknown diagram_type: {diagram_type!r}")

        return {"ok": True, "output_path": output_path, "error": None}
    except Exception as exc:
        return {"ok": False, "output_path": output_path, "error": str(exc)}


@mcp.tool()
def render_formula_png(
    formula_tex: str,
    output_path: str,
    fig_width: float = 7.0,
    fig_height: float = 1.0,
    fontsize: int = 22,
    color: str = "#1F497D",
    background: str | None = None,
    dpi: int = 160,
) -> dict[str, Any]:
    """Render a LaTeX/mathtext formula as a transparent PNG using matplotlib mathtext.

    No LaTeX installation required. Uses a single $...$ mathtext block.
    Avoid adjacent $A$ $B$ pairs and \\& (use \\mathrm{and} instead).
    Returns: ok, output_path, width_px, height_px, error.
    """
    from cumt_ppt_mcp.ppt_utils import render_formula_png as _impl
    return _impl(formula_tex, output_path, fig_width, fig_height, fontsize, color, background, dpi)


@mcp.tool()
def add_breadcrumb_strip(
    pptx_path: str,
    output_path: str,
    sections: list[dict[str, str]],
    slide_section_map: dict[str, int | None],
    top_inch: float = 1.0,
    height_inch: float = 0.32,
    active_color: str = "#FFC000",
    inactive_color: str = "#2E5FA3",
) -> dict[str, Any]:
    """Add a section-breadcrumb navigation strip to every slide in a PPTX file.

    sections: list of {label, name} dicts, e.g. [{"label":"01","name":"研究背景"}, ...].
    slide_section_map: maps 1-based slide number (as string) to 0-based section index or null.
    Saves the modified deck to output_path.
    Returns: ok, slide_count, error.
    """
    try:
        from pptx import Presentation
        from cumt_ppt_mcp.ppt_utils import add_breadcrumb_strip as _impl

        section_tuples = [(s["label"], s["name"]) for s in sections]
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            slide_num_str = str(i + 1)
            active_idx = slide_section_map.get(slide_num_str)
            _impl(slide, section_tuples, active_idx,
                  top_inch=top_inch, height_inch=height_inch,
                  active_color=active_color, inactive_color=inactive_color)
        prs.save(output_path)
        return {"ok": True, "slide_count": len(prs.slides), "output_path": output_path, "error": None}
    except Exception as exc:
        return {"ok": False, "slide_count": 0, "output_path": output_path, "error": str(exc)}


def main() -> None:
    mcp.run()


if __name__ == "__main__":
    main()
