from __future__ import annotations

import math
import os
import re
import shutil
import subprocess
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt


EMU_PER_INCH = 914400
PLACEHOLDER_WORDS = [
    "单击此处",
    "添加标题",
    "添加文本",
    "Key Words",
    "Question",
    "Lorem",
    "Vivamus",
    "Your Text",
    "INSERT",
    "TODO",
]

# Keep this second assignment in UTF-8 so placeholder checks work in Chinese decks.
PLACEHOLDER_WORDS = [
    "单击此处",
    "添加标题",
    "添加文本",
    "Key Words",
    "Question",
    "Lorem",
    "Vivamus",
    "Your Text",
    "INSERT",
    "TODO",
]

PROJECT_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_LOGO_PATH = PROJECT_ROOT / "assets" / "cumt_logo.png"


class PptToolError(ValueError):
    """Clear, user-facing PPT tool error."""


def _path(value: str | os.PathLike[str]) -> Path:
    return Path(value).expanduser().resolve()


def validate_pptx_path(pptx_path: str | os.PathLike[str]) -> Path:
    path = _path(pptx_path)
    if not path.exists():
        raise PptToolError(f"PPTX file does not exist: {path}")
    if not path.is_file():
        raise PptToolError(f"Path is not a file: {path}")
    if path.suffix.lower() != ".pptx":
        raise PptToolError(f"Expected a .pptx file: {path}")
    return path


def validate_output_path(output_path: str | os.PathLike[str], input_path: Path | None = None) -> Path:
    path = _path(output_path)
    if path.suffix.lower() != ".pptx":
        raise PptToolError(f"Output path must end with .pptx: {path}")
    if input_path and path == input_path.resolve():
        raise PptToolError("Output path must be different from input PPTX path.")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def emu_to_inches(value: int | Emu) -> float:
    return round(int(value) / EMU_PER_INCH, 3)


def inches_to_emu(value: float | int) -> Emu:
    return Inches(float(value))


def shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
    }


def has_text(shape: Any) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and shape.text_frame and shape.text.strip())


def is_picture(shape: Any) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def is_table(shape: Any) -> bool:
    return bool(getattr(shape, "has_table", False))


def iter_text_runs(shape: Any):
    if not has_text(shape):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            yield paragraph, run


def get_run_font_info(run: Any) -> dict[str, Any]:
    font = run.font
    return {
        "text": run.text,
        "font_name": font.name,
        "size_pt": round(font.size.pt, 1) if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
    }


def shape_font_summary(shape: Any) -> list[dict[str, Any]]:
    fonts = []
    if not has_text(shape):
        return fonts
    for _paragraph, run in iter_text_runs(shape):
        if run.text:
            fonts.append(get_run_font_info(run))
    return fonts


def text_shape_info(shape: Any) -> dict[str, Any]:
    return {
        "text": shape.text.strip(),
        "bounds_in": shape_bounds(shape),
        "fonts": shape_font_summary(shape),
    }


def _font_size_max(shape: Any) -> float:
    sizes = []
    for _paragraph, run in iter_text_runs(shape) or []:
        if run.font.size:
            sizes.append(run.font.size.pt)
    return max(sizes) if sizes else 0.0


def _is_title_shape(shape: Any, slide_height: int) -> bool:
    if not has_text(shape):
        return False
    top_ratio = int(shape.top) / max(slide_height, 1)
    text = shape.text.strip()
    if top_ratio <= 0.16 and len(text) <= 80:
        return True
    return _font_size_max(shape) >= 24 and top_ratio <= 0.28


def detect_title(slide: Any, slide_height: int) -> str:
    candidates = [shape for shape in slide.shapes if _is_title_shape(shape, slide_height)]
    if not candidates:
        candidates = [shape for shape in slide.shapes if has_text(shape)]
    if not candidates:
        return ""
    candidates.sort(key=lambda shape: (shape.top, -_font_size_max(shape), shape.left))
    return candidates[0].text.strip().replace("\n", " ")


def inspect_ppt(pptx_path: str) -> dict[str, Any]:
    path = validate_pptx_path(pptx_path)
    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        slides.append(
            {
                "slide_index": idx,
                "title": detect_title(slide, prs.slide_height),
                "text_box_count": sum(1 for shape in slide.shapes if has_text(shape)),
                "image_count": sum(1 for shape in slide.shapes if is_picture(shape)),
                "table_count": sum(1 for shape in slide.shapes if is_table(shape)),
            }
        )
    return {
        "pptx_path": str(path),
        "slide_count": len(prs.slides),
        "page_size": {
            "width_in": emu_to_inches(prs.slide_width),
            "height_in": emu_to_inches(prs.slide_height),
            "width_emu": int(prs.slide_width),
            "height_emu": int(prs.slide_height),
        },
        "slides": slides,
    }


def _dense_hints(slide: Any, prs: Presentation) -> list[str]:
    hints: list[str] = []
    text_shapes = [shape for shape in slide.shapes if has_text(shape)]
    total_chars = sum(len(shape.text.strip()) for shape in text_shapes)
    pictures = [shape for shape in slide.shapes if is_picture(shape)]
    small_fonts = []
    out_of_bounds = []
    for shape in slide.shapes:
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
            out_of_bounds.append(shape.name)
        for _paragraph, run in iter_text_runs(shape) or []:
            if run.font.size and run.font.size.pt < 10:
                small_fonts.append({"shape": shape.name, "text": run.text[:30], "size_pt": round(run.font.size.pt, 1)})
    if len(text_shapes) >= 12:
        hints.append(f"High number of text boxes: {len(text_shapes)}")
    if total_chars >= 900:
        hints.append(f"Possible dense slide text: {total_chars} characters")
    if len(pictures) >= 6:
        hints.append(f"Many pictures on one slide: {len(pictures)}")
    if small_fonts:
        hints.append(f"Small font risk: {len(small_fonts)} runs under 10 pt")
    if out_of_bounds:
        hints.append(f"Shapes may be outside slide bounds: {', '.join(out_of_bounds[:5])}")
    return hints


def inspect_slide(pptx_path: str, slide_index: int) -> dict[str, Any]:
    path = validate_pptx_path(pptx_path)
    prs = Presentation(str(path))
    if slide_index < 1 or slide_index > len(prs.slides):
        raise PptToolError(f"slide_index must be between 1 and {len(prs.slides)}.")
    slide = prs.slides[slide_index - 1]
    text_boxes = []
    pictures = []
    tables = []
    for shape in slide.shapes:
        if has_text(shape):
            text_boxes.append(text_shape_info(shape))
        if is_picture(shape):
            pictures.append({"name": shape.name, "bounds_in": shape_bounds(shape)})
        if is_table(shape):
            tables.append(
                {
                    "name": shape.name,
                    "bounds_in": shape_bounds(shape),
                    "rows": len(shape.table.rows),
                    "columns": len(shape.table.columns),
                }
            )
    return {
        "pptx_path": str(path),
        "slide_index": slide_index,
        "title": detect_title(slide, prs.slide_height),
        "text_boxes": text_boxes,
        "pictures": pictures,
        "tables": tables,
        "hints": _dense_hints(slide, prs),
    }


def _ensure_rfonts(run: Any, latin_font: str, east_asia_font: str) -> None:
    r_pr = run._r.get_or_add_rPr()
    for tag, face in (("latin", latin_font), ("ea", east_asia_font), ("cs", latin_font)):
        element = r_pr.find(qn(f"a:{tag}"))
        if element is None:
            element = OxmlElement(f"a:{tag}")
            r_pr.append(element)
        element.set("typeface", face)
    run.font.name = latin_font


def apply_font_rules(
    pptx_path: str,
    output_path: str,
    title_font: str = "黑体",
    body_font: str = "宋体",
    latin_font: str = "Times New Roman",
) -> dict[str, Any]:
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    prs = Presentation(str(input_path))
    changed_runs = 0
    title_shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not has_text(shape):
                continue
            is_title = _is_title_shape(shape, prs.slide_height)
            if is_title:
                title_shapes += 1
            east_font = title_font if is_title else body_font
            for _paragraph, run in iter_text_runs(shape) or []:
                _ensure_rfonts(run, latin_font, east_font)
                if is_title:
                    run.font.bold = True
                changed_runs += 1
    prs.save(str(out_path))
    return {
        "ok": True,
        "input_path": str(input_path),
        "output_path": str(out_path),
        "changed_runs": changed_runs,
        "title_shapes": title_shapes,
        "title_font": title_font,
        "body_font": body_font,
        "latin_font": latin_font,
    }


def _logo_candidate_evaluation(shape: Any, slide_width: int, slide_height: int) -> dict[str, Any]:
    """Return a strict, explainable CUMT header-logo candidate evaluation.

    The rules intentionally reject ambiguous images. Template decks often contain
    thesis figures, screenshots, and diagrams near the upper half of a slide; only
    small, top-right, header-like pictures should be reused as logo placement.
    """
    bounds = shape_bounds(shape)
    width_ratio = int(shape.width) / max(int(slide_width), 1)
    height_ratio = int(shape.height) / max(int(slide_height), 1)
    area_ratio = width_ratio * height_ratio
    left_ratio = int(shape.left) / max(int(slide_width), 1)
    top_ratio = int(shape.top) / max(int(slide_height), 1)
    right_ratio = int(shape.left + shape.width) / max(int(slide_width), 1)
    bottom_ratio = int(shape.top + shape.height) / max(int(slide_height), 1)
    aspect_ratio = int(shape.width) / max(int(shape.height), 1)
    reasons: list[str] = []

    # Hard forbid large pictures first: figures, screenshots, background images,
    # and half-slide visuals must never become header logos.
    if width_ratio > 0.20:
        reasons.append("width exceeds 20% of slide")
    if height_ratio > 0.20:
        reasons.append("height exceeds 20% of slide")
    if area_ratio > 0.05:
        reasons.append("area exceeds 5% of slide")

    # Reusable CUMT logo placement should be in the page header, not the body.
    if not (left_ratio >= 0.68 or right_ratio >= 0.90):
        reasons.append("not close enough to right edge")
    if not (top_ratio <= 0.13 and bottom_ratio <= 0.24):
        reasons.append("not in top header area")

    # Preferred strict size: small enough to be page chrome rather than content.
    if width_ratio > 0.15:
        reasons.append("width exceeds 15% header-logo limit")
    if height_ratio > 0.15:
        reasons.append("height exceeds 15% header-logo limit")
    if area_ratio > 0.03:
        reasons.append("area exceeds 3% header-logo limit")

    # CUMT logo+wordmark is usually horizontal; allow circular marks but reject
    # extreme banners and tall/skinny diagrams.
    if not (0.6 <= aspect_ratio <= 8.5):
        reasons.append("aspect ratio outside header-logo range")

    # A picture starting deep inside the body is probably a figure even if its
    # right edge reaches the page edge.
    if top_ratio > 0.16 or left_ratio < 0.55:
        reasons.append("located in body/content region")

    return {
        "is_candidate": len(reasons) == 0,
        "bounds_in": bounds,
        "width_ratio": round(width_ratio, 4),
        "height_ratio": round(height_ratio, 4),
        "area_ratio": round(area_ratio, 4),
        "aspect_ratio": round(aspect_ratio, 3),
        "reasons": reasons,
    }


def is_logo_candidate(shape: Any, slide_width: int, slide_height: int) -> bool:
    """Strictly decide whether a picture can be treated as a header logo."""
    return bool(is_picture(shape) and _logo_candidate_evaluation(shape, slide_width, slide_height)["is_candidate"])


def _classify_picture_shape(shape: Any, prs: Presentation) -> dict[str, Any]:
    eval_result = _logo_candidate_evaluation(shape, int(prs.slide_width), int(prs.slide_height))
    bounds = eval_result["bounds_in"]
    top_ratio = int(shape.top) / max(int(prs.slide_height), 1)
    area_ratio = eval_result["area_ratio"]
    width_ratio = eval_result["width_ratio"]
    height_ratio = eval_result["height_ratio"]
    right_ratio = int(shape.left + shape.width) / max(int(prs.slide_width), 1)
    bottom_ratio = int(shape.top + shape.height) / max(int(prs.slide_height), 1)

    if eval_result["is_candidate"]:
        classification = "header_logo"
        reusable = True
        reason = "small repeated-style picture in top-right header area"
    elif area_ratio > 0.05 or width_ratio > 0.20 or height_ratio > 0.20:
        classification = "content_image"
        reusable = False
        reason = "area or dimensions too large for a header logo; treat as thesis/content figure"
    elif top_ratio >= 0.72 and area_ratio <= 0.04:
        classification = "background_decoration"
        reusable = False
        reason = "small lower-page decoration; not a header logo"
    elif right_ratio >= 0.85 and bottom_ratio <= 0.30:
        classification = "unknown_image"
        reusable = False
        reason = "near header but failed logo constraints: " + "; ".join(eval_result["reasons"])
    else:
        classification = "content_image" if top_ratio > 0.18 else "unknown_image"
        reusable = False
        reason = "not in right-top header logo region"
    return {
        "bounds_in": bounds,
        "width_ratio": eval_result["width_ratio"],
        "height_ratio": eval_result["height_ratio"],
        "area_ratio": eval_result["area_ratio"],
        "aspect_ratio": eval_result["aspect_ratio"],
        "classification": classification,
        "is_template_feature": reusable,
        "reason": reason,
        "failed_logo_rules": eval_result["reasons"],
    }


def find_header_logo(slide: Any, prs: Presentation) -> Any | None:
    candidates = [shape for shape in slide.shapes if is_logo_candidate(shape, int(prs.slide_width), int(prs.slide_height))]
    candidates.sort(key=lambda s: (_logo_candidate_evaluation(s, int(prs.slide_width), int(prs.slide_height))["area_ratio"], s.top, -s.left))
    return candidates[0] if candidates else None


def _right_top_pictures(slide: Any, prs: Presentation) -> list[Any]:
    """Backward-compatible strict header-logo finder."""
    pics = [shape for shape in slide.shapes if is_logo_candidate(shape, int(prs.slide_width), int(prs.slide_height))]
    pics.sort(key=lambda s: (_logo_candidate_evaluation(s, int(prs.slide_width), int(prs.slide_height))["area_ratio"], s.top, -s.left))
    return pics


def extract_logo_style(template_pptx_path: str) -> dict[str, Any]:
    """Extract reusable header-logo style only; never returns content pictures."""
    path = validate_pptx_path(template_pptx_path)
    prs = Presentation(str(path))
    logos = []
    for idx, slide in enumerate(prs.slides, start=1):
        logo = find_header_logo(slide, prs)
        if logo is not None:
            logos.append({"slide_index": idx, "bounds_in": shape_bounds(logo)})
    return {"ok": True, "template_pptx_path": str(path), "header_logos": logos, "default_logo": logos[0] if logos else None}


def normalize_logo(
    pptx_path: str,
    output_path: str,
    reference_slide_index: int = 3,
    start_slide_index: int = 4,
    skip_slide_indices: list[int] | None = None,
    logo_asset_path: str | None = None,
) -> dict[str, Any]:
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    logo_path = _logo_asset_path(logo_asset_path)
    prs = Presentation(str(input_path))
    if reference_slide_index < 1 or reference_slide_index > len(prs.slides):
        raise PptToolError(f"reference_slide_index must be between 1 and {len(prs.slides)}.")
    ref_slide = prs.slides[reference_slide_index - 1]
    ref_pics = _right_top_pictures(ref_slide, prs)
    if not ref_pics:
        ref_bounds = {
            "left": emu_to_inches(prs.slide_width) - 2.15,
            "top": 0.1,
            "width": 1.7,
            "height": 0.42,
        }
    else:
        ref = ref_pics[0]
        ref_bounds = shape_bounds(ref)
    skip = set(skip_slide_indices or [])
    modified: list[int] = []
    not_found: list[int] = []
    for idx in range(max(start_slide_index, 1), len(prs.slides) + 1):
        if idx == reference_slide_index or idx in skip:
            continue
        slide = prs.slides[idx - 1]
        pics = _right_top_pictures(slide, prs)
        if pics:
            for old in pics:
                _delete_shape(old)
        else:
            not_found.append(idx)
        slide.shapes.add_picture(
            str(logo_path),
            inches_to_emu(ref_bounds["left"]),
            inches_to_emu(ref_bounds["top"]),
            width=inches_to_emu(ref_bounds["width"]),
        )
        modified.append(idx)
    prs.save(str(out_path))
    return {
        "ok": True,
        "input_path": str(input_path),
        "output_path": str(out_path),
        "logo_asset_path": str(logo_path),
        "reference_slide_index": reference_slide_index,
        "reference_logo": ref_bounds,
        "modified_slide_indices": modified,
        "not_found_slide_indices": not_found,
    }


def _selected_slide_indices(total: int, slide_indices: list[int] | None) -> list[int]:
    if not slide_indices:
        return list(range(1, total + 1))
    bad = [idx for idx in slide_indices if idx < 1 or idx > total]
    if bad:
        raise PptToolError(f"slide_indices out of range 1..{total}: {bad}")
    return slide_indices


def _scale_shape_keep_ratio(shape: Any, max_w: int, max_h: int) -> bool:
    ratio = min(max_w / shape.width, max_h / shape.height, 1.0)
    if ratio >= 0.999:
        return False
    shape.width = int(shape.width * ratio)
    shape.height = int(shape.height * ratio)
    return True


def resize_images_keep_ratio(
    pptx_path: str,
    output_path: str,
    slide_indices: list[int] | None = None,
    mode: str = "fit_content_area",
    max_width_ratio: float = 0.45,
    max_height_ratio: float = 0.58,
) -> dict[str, Any]:
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    if mode not in {"center", "fit_content_area", "scale"}:
        raise PptToolError("mode must be one of: center, fit_content_area, scale")
    if not 0.05 <= max_width_ratio <= 1.0 or not 0.05 <= max_height_ratio <= 1.0:
        raise PptToolError("max_width_ratio and max_height_ratio must be between 0.05 and 1.0")
    prs = Presentation(str(input_path))
    selected = _selected_slide_indices(len(prs.slides), slide_indices)
    modified: list[dict[str, Any]] = []
    content_top = int(prs.slide_height * 0.16)
    for idx in selected:
        slide = prs.slides[idx - 1]
        pictures = [shape for shape in slide.shapes if is_picture(shape)]
        for pic_shape in pictures:
            old = shape_bounds(pic_shape)
            max_w = int(prs.slide_width * max_width_ratio)
            max_h = int(prs.slide_height * max_height_ratio)
            changed = _scale_shape_keep_ratio(pic_shape, max_w, max_h)
            if mode in {"center", "fit_content_area"}:
                pic_shape.left = int((prs.slide_width - pic_shape.width) / 2)
                if mode == "fit_content_area":
                    pic_shape.top = int(content_top + (prs.slide_height - content_top - pic_shape.height) / 2)
                else:
                    pic_shape.top = int((prs.slide_height - pic_shape.height) / 2)
                changed = True
            elif mode == "scale" and pic_shape.top < content_top:
                pic_shape.top = content_top
                changed = True
            if changed:
                modified.append({"slide_index": idx, "name": pic_shape.name, "old_bounds_in": old, "new_bounds_in": shape_bounds(pic_shape)})
    prs.save(str(out_path))
    return {
        "ok": True,
        "input_path": str(input_path),
        "output_path": str(out_path),
        "mode": mode,
        "modified_images": modified,
    }


def _split_font_runs(text: str) -> list[tuple[str, bool]]:
    parts = re.findall(r"[A-Za-z0-9 .,%+\-_/():=<>×^]+|[^\x00-\x7f]+|.", str(text))
    return [(part, bool(re.fullmatch(r"[A-Za-z0-9 .,%+\-_/():=<>×^]+", part))) for part in parts if part]


def _add_text(slide: Any, left: int, top: int, width: int, height: int, text: str, size_pt: float, bold: bool = False, align: Any = PP_ALIGN.CENTER) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for part, is_latin in _split_font_runs(text):
        run = paragraph.add_run()
        run.text = part
        _ensure_rfonts(run, "Times New Roman", "宋体")
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(30, 30, 30)
        if is_latin:
            run.font.name = "Times New Roman"


def _add_line(slide: Any, left: int, top: int, width: int, color: RGBColor = RGBColor(29, 67, 113), pt: float = 1.5) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color


def make_three_line_table(
    pptx_path: str,
    output_path: str,
    slide_index: int,
    table_data: list[list[Any]],
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    height: float | None = None,
) -> dict[str, Any]:
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    if not table_data or not all(isinstance(row, list) and row for row in table_data):
        raise PptToolError("table_data must be a non-empty 2D array.")
    columns = len(table_data[0])
    if any(len(row) != columns for row in table_data):
        raise PptToolError("All rows in table_data must have the same number of columns.")
    prs = Presentation(str(input_path))
    if slide_index < 1 or slide_index > len(prs.slides):
        raise PptToolError(f"slide_index must be between 1 and {len(prs.slides)}.")
    slide = prs.slides[slide_index - 1]
    left_emu = inches_to_emu(left if left is not None else 1.0)
    top_emu = inches_to_emu(top if top is not None else 1.35)
    width_emu = inches_to_emu(width if width is not None else 11.3)
    height_emu = inches_to_emu(height if height is not None else 4.6)
    rows = len(table_data)
    row_h = int(height_emu / rows)
    col_w = int(width_emu / columns)
    _add_line(slide, left_emu, top_emu, width_emu, pt=1.6)
    _add_line(slide, left_emu, top_emu + row_h, width_emu, pt=1.2)
    _add_line(slide, left_emu, top_emu + height_emu, width_emu, pt=1.6)
    for r_idx, row in enumerate(table_data):
        for c_idx, value in enumerate(row):
            size = 15 if r_idx == 0 else 13
            bold = r_idx == 0
            _add_text(
                slide,
                left_emu + c_idx * col_w,
                top_emu + r_idx * row_h,
                col_w,
                row_h,
                str(value),
                size,
                bold,
                PP_ALIGN.CENTER,
            )
    prs.save(str(out_path))
    return {
        "ok": True,
        "input_path": str(input_path),
        "output_path": str(out_path),
        "slide_index": slide_index,
        "rows": rows,
        "columns": columns,
        "bounds_in": {
            "left": emu_to_inches(left_emu),
            "top": emu_to_inches(top_emu),
            "width": emu_to_inches(width_emu),
            "height": emu_to_inches(height_emu),
        },
    }


def export_preview(pptx_path: str, output_dir: str) -> dict[str, Any]:
    input_path = validate_pptx_path(pptx_path)
    out_dir = _path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    errors: list[str] = []
    def preview_files() -> list[str]:
        paths = {path.resolve() for pattern in ("*.PNG", "*.png") for path in out_dir.glob(pattern)}
        return [str(path) for path in sorted(paths, key=lambda item: item.name.lower())]

    try:
        import win32com.client  # type: ignore

        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = True
        pres = app.Presentations.Open(str(input_path), False, False, False)
        pres.Export(str(out_dir), "PNG", 1600, 900)
        slide_count = pres.Slides.Count
        pres.Close()
        app.Quit()
        files = preview_files()
        return {
            "ok": True,
            "method": "powerpoint_com",
            "pptx_path": str(input_path),
            "output_dir": str(out_dir),
            "slide_count": slide_count,
            "preview_files": files,
        }
    except Exception as exc:  # pragma: no cover - environment-specific
        errors.append(f"PowerPoint COM export failed: {exc}")
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "png", "--outdir", str(out_dir), str(input_path)],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            if result.returncode == 0:
                files = preview_files()
                return {
                    "ok": True,
                    "method": "libreoffice",
                    "pptx_path": str(input_path),
                    "output_dir": str(out_dir),
                    "preview_files": files,
                    "stdout": result.stdout,
                }
            errors.append(f"LibreOffice returned {result.returncode}: {result.stderr or result.stdout}")
        except Exception as exc:  # pragma: no cover - environment-specific
            errors.append(f"LibreOffice export failed: {exc}")
    else:
        errors.append("LibreOffice/soffice executable was not found.")
    return {
        "ok": False,
        "pptx_path": str(input_path),
        "output_dir": str(out_dir),
        "errors": errors,
        "suggestion": "Open the PPTX in PowerPoint and export slides as PNG, or install LibreOffice/PowerPoint automation support.",
    }


def _picture_aspect_risk(shape: Any) -> dict[str, Any] | None:
    try:
        blob = shape.image.blob
        with Image.open(BytesIO(blob)) as img:
            image_ratio = img.width / img.height
        shape_ratio = shape.width / shape.height
        if abs(image_ratio - shape_ratio) / image_ratio > 0.08:
            return {
                "name": shape.name,
                "bounds_in": shape_bounds(shape),
                "image_aspect": round(image_ratio, 3),
                "shape_aspect": round(shape_ratio, 3),
            }
    except Exception:
        return None
    return None


def _unique_report_path(pptx_path: Path) -> Path:
    base = pptx_path.with_name(pptx_path.stem + "_quality_report.md")
    if not base.exists():
        return base
    for i in range(1, 100):
        candidate = pptx_path.with_name(f"{pptx_path.stem}_quality_report_{i}.md")
        if not candidate.exists():
            return candidate
    raise PptToolError("Could not create a unique report path.")


def check_ppt_quality(
    pptx_path: str,
    expected_title: str | None = None,
    required_numbers: list[str] | None = None,
    report_path: str | None = None,
) -> dict[str, Any]:
    path = validate_pptx_path(pptx_path)
    prs = Presentation(str(path))
    all_text = []
    small_fonts = []
    aspect_risks = []
    dense_slides = []
    logo_candidates = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if has_text(shape):
                all_text.append(shape.text)
                for _paragraph, run in iter_text_runs(shape) or []:
                    if run.font.size and run.font.size.pt < 10:
                        small_fonts.append({"slide_index": idx, "text": run.text[:40], "size_pt": round(run.font.size.pt, 1)})
            if is_picture(shape):
                risk = _picture_aspect_risk(shape)
                if risk:
                    risk["slide_index"] = idx
                    aspect_risks.append(risk)
        hints = _dense_hints(slide, prs)
        if hints:
            dense_slides.append({"slide_index": idx, "hints": hints})
        pics = _right_top_pictures(slide, prs)
        if pics:
            logo_candidates.append((idx, pics[0]))
    text = "\n".join(all_text)
    required_numbers = required_numbers or []
    missing_numbers = [num for num in required_numbers if num not in text]
    placeholders = [word for word in PLACEHOLDER_WORDS if word in text]
    expected_title_present = None if expected_title is None else expected_title in text
    logo_inconsistency = []
    if len(logo_candidates) >= 2:
        ref_idx, ref = logo_candidates[0]
        for idx, shape in logo_candidates[1:]:
            if (
                abs(shape.width - ref.width) > EMU_PER_INCH * 0.03
                or abs(shape.height - ref.height) > EMU_PER_INCH * 0.03
                or abs(shape.left - ref.left) > EMU_PER_INCH * 0.04
                or abs(shape.top - ref.top) > EMU_PER_INCH * 0.04
            ):
                logo_inconsistency.append({"slide_index": idx, "reference_slide_index": ref_idx, "bounds_in": shape_bounds(shape)})
    report_file = _path(report_path) if report_path else _unique_report_path(path)
    if report_file.suffix.lower() != ".md":
        raise PptToolError(f"Quality report path must end with .md: {report_file}")
    report_file.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        "# PPT Quality Report",
        "",
        f"- PPTX: `{path}`",
        f"- Slides: {len(prs.slides)}",
        f"- Expected title present: {expected_title_present}",
        f"- Missing required numbers: {missing_numbers}",
        f"- Placeholder words found: {placeholders}",
        f"- Small font runs under 10 pt: {len(small_fonts)}",
        f"- Picture aspect risks: {len(aspect_risks)}",
        f"- Logo consistency risks: {len(logo_inconsistency)}",
        "",
        "## Dense Slide Hints",
    ]
    for item in dense_slides:
        lines.append(f"- Slide {item['slide_index']}: {'; '.join(item['hints'])}")
    lines.extend(["", "## Small Fonts"])
    for item in small_fonts[:50]:
        lines.append(f"- Slide {item['slide_index']}: {item['size_pt']} pt, `{item['text']}`")
    lines.extend(["", "## Picture Aspect Risks"])
    for item in aspect_risks[:50]:
        lines.append(f"- Slide {item['slide_index']}: {item['name']} aspect {item['shape_aspect']} vs original {item['image_aspect']}")
    lines.extend(["", "## Logo Risks"])
    for item in logo_inconsistency[:50]:
        lines.append(f"- Slide {item['slide_index']}: differs from slide {item['reference_slide_index']}")
    report_file.write_text("\n".join(lines), encoding="utf-8")
    return {
        "ok": True,
        "pptx_path": str(path),
        "expected_title_present": expected_title_present,
        "missing_required_numbers": missing_numbers,
        "placeholder_words_found": placeholders,
        "small_font_count": len(small_fonts),
        "picture_aspect_risk_count": len(aspect_risks),
        "logo_inconsistency_count": len(logo_inconsistency),
        "dense_slide_hint_count": len(dense_slides),
        "report_path": str(report_file),
    }


def validate_pdf_path(pdf_path: str | os.PathLike[str]) -> Path:
    path = _path(pdf_path)
    if not path.exists():
        raise PptToolError(f"PDF file does not exist: {path}")
    if not path.is_file():
        raise PptToolError(f"Path is not a file: {path}")
    if path.suffix.lower() != ".pdf":
        raise PptToolError(f"Expected a .pdf file: {path}")
    return path


def _rgb_tuple(rgb: RGBColor | None) -> tuple[int, int, int] | None:
    if rgb is None:
        return None
    return (int(rgb[0]), int(rgb[1]), int(rgb[2]))


def _shape_fill_rgb(shape: Any) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        color = fill.fore_color
        if color.rgb is not None:
            return _rgb_tuple(color.rgb)
    except Exception:
        pass
    try:
        srgb = shape._element.xpath(".//a:solidFill/a:srgbClr")
        if srgb:
            value = srgb[0].get("val")
            if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
                return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return None


def _text_color_rgb(run: Any) -> tuple[int, int, int] | None:
    try:
        rgb = run.font.color.rgb
        return _rgb_tuple(rgb) if rgb is not None else None
    except Exception:
        return None


def _most_common(items: Iterable[Any], limit: int = 5) -> list[dict[str, Any]]:
    counts: dict[Any, int] = {}
    for item in items:
        if item is None:
            continue
        counts[item] = counts.get(item, 0) + 1
    return [{"value": item, "count": count} for item, count in sorted(counts.items(), key=lambda kv: kv[1], reverse=True)[:limit]]


def _delete_shape(shape: Any) -> None:
    shape._element.getparent().remove(shape._element)


def _send_to_back(shape: Any) -> None:
    tree = shape._element.getparent()
    tree.remove(shape._element)
    tree.insert(2, shape._element)


def _logo_asset_path(logo_asset_path: str | None = None) -> Path:
    path = _path(logo_asset_path) if logo_asset_path else DEFAULT_LOGO_PATH
    if not path.exists():
        raise PptToolError(f"CUMT logo asset does not exist: {path}")
    if path.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
        raise PptToolError(f"CUMT logo asset must be an image file: {path}")
    return path


def check_cumt_logo_asset(logo_asset_path: str | None = None) -> dict[str, Any]:
    """Check the reusable CUMT logo asset for black background, alpha, and size risks."""
    path = _logo_asset_path(logo_asset_path)
    with Image.open(path) as image:
        rgba = image.convert("RGBA")
        rgb = rgba.convert("RGB")
        w, h = rgba.size
        pixels = list(rgba.getdata())
        alpha_values = [p[3] for p in pixels]
        transparent_ratio = sum(1 for a in alpha_values if a < 250) / max(len(alpha_values), 1)
        rgb_pixels = list(rgb.getdata())
        dark_ratio = sum(1 for r, g, b in rgb_pixels if r < 35 and g < 35 and b < 35) / max(len(rgb_pixels), 1)
        border = []
        for x in range(w):
            border.append(rgb.getpixel((x, 0)))
            border.append(rgb.getpixel((x, h - 1)))
        for y in range(h):
            border.append(rgb.getpixel((0, y)))
            border.append(rgb.getpixel((w - 1, y)))
        black_border_ratio = sum(1 for r, g, b in border if r < 35 and g < 35 and b < 35) / max(len(border), 1)
        warnings = []
        if w < 200 or h < 50:
            warnings.append("Logo image resolution may be too small.")
        if transparent_ratio > 0.05:
            warnings.append("Logo contains notable transparency; compositing on white is recommended.")
        if dark_ratio > 0.65 or black_border_ratio > 0.65:
            warnings.append("Logo may have a black background or black border.")
        return {
            "ok": True,
            "logo_asset_path": str(path),
            "size_px": {"width": w, "height": h},
            "mode": image.mode,
            "transparent_ratio": round(transparent_ratio, 4),
            "dark_pixel_ratio": round(dark_ratio, 4),
            "black_border_ratio": round(black_border_ratio, 4),
            "warnings": warnings,
        }


def add_cumt_logo(
    pptx_path: str,
    output_path: str,
    slide_indices: list[int] | None = None,
    logo_asset_path: str | None = None,
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    skip_slide_indices: list[int] | None = None,
) -> dict[str, Any]:
    """Add the fixed CUMT logo asset to selected slides and save a new PPTX."""
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    logo_path = _logo_asset_path(logo_asset_path)
    prs = Presentation(str(input_path))
    selected = _selected_slide_indices(len(prs.slides), slide_indices)
    skip = set(skip_slide_indices or [])
    logo_w = width if width is not None else 1.7
    logo_left = left if left is not None else emu_to_inches(prs.slide_width) - logo_w - 0.45
    logo_top = top if top is not None else 0.1
    modified = []
    for idx in selected:
        if idx in skip:
            continue
        slide = prs.slides[idx - 1]
        for old in _right_top_pictures(slide, prs):
            _delete_shape(old)
        shape = slide.shapes.add_picture(str(logo_path), inches_to_emu(logo_left), inches_to_emu(logo_top), width=inches_to_emu(logo_w))
        modified.append({"slide_index": idx, "bounds_in": shape_bounds(shape)})
    prs.save(str(out_path))
    return {
        "ok": True,
        "input_path": str(input_path),
        "output_path": str(out_path),
        "logo_asset_path": str(logo_path),
        "modified_slides": modified,
    }


def _template_layout_type(slide: Any, prs: Presentation, slide_index: int) -> str:
    title = detect_title(slide, prs.slide_height)
    text_count = sum(1 for shape in slide.shapes if has_text(shape))
    pic_count = sum(1 for shape in slide.shapes if is_picture(shape))
    if slide_index == 1:
        return "cover"
    if "目录" in title or "Contents" in title:
        return "contents"
    if text_count <= 3 and pic_count <= 1 and any(len(shape.text.strip()) <= 12 for shape in slide.shapes if has_text(shape)):
        return "section"
    if "结论" in title or "展望" in title:
        return "conclusion"
    if "感谢" in title or "致谢" in title or "批评指正" in title:
        return "thanks"
    if pic_count >= 2:
        return "figure"
    if text_count >= 10:
        return "multi_block"
    return "left_text_right_figure" if pic_count == 1 else "content"


def inspect_template_style(template_pptx_path: str) -> dict[str, Any]:
    """Inspect a reference template deck and summarize reusable style features."""
    path = validate_pptx_path(template_pptx_path)
    prs = Presentation(str(path))
    backgrounds = []
    fills = []
    title_fonts = []
    body_fonts = []
    body_sizes = []
    title_sizes = []
    title_colors = []
    body_colors = []
    emphasis_colors = []
    header_candidates = []
    logo_positions = []
    image_classifications = []
    page_numbers = []
    layouts = []
    for idx, slide in enumerate(prs.slides, start=1):
        # Background colors from full-slide or top-band rectangles.
        picture_no = 0
        for shape in slide.shapes:
            color = _shape_fill_rgb(shape)
            if color:
                fills.append(color)
            if (
                color
                and shape.left <= Inches(0.1)
                and shape.top <= Inches(0.25)
                and shape.width >= prs.slide_width * 0.7
                and 0.15 <= emu_to_inches(shape.height) <= 1.4
            ):
                header_candidates.append({"slide_index": idx, "bounds_in": shape_bounds(shape), "color": color})
            if has_text(shape):
                is_title = _is_title_shape(shape, prs.slide_height)
                for _paragraph, run in iter_text_runs(shape) or []:
                    if not run.text.strip():
                        continue
                    size = round(run.font.size.pt, 1) if run.font.size else None
                    color_text = _text_color_rgb(run)
                    if color_text and color_text[0] > 150 and color_text[1] < 80 and color_text[2] < 80:
                        emphasis_colors.append(color_text)
                    if is_title:
                        title_fonts.append(run.font.name)
                        title_sizes.append(size)
                        title_colors.append(color_text)
                    else:
                        body_fonts.append(run.font.name)
                        body_sizes.append(size)
                        body_colors.append(color_text)
            if is_picture(shape):
                picture_no += 1
                classification = _classify_picture_shape(shape, prs)
                classification.update({"slide_index": idx, "picture_index": picture_no, "name": getattr(shape, "name", "")})
                image_classifications.append(classification)
                if classification["classification"] == "header_logo":
                    logo_positions.append({"slide_index": idx, "picture_index": picture_no, "bounds_in": classification["bounds_in"]})
        title = detect_title(slide, prs.slide_height)
        title_shapes = [shape for shape in slide.shapes if _is_title_shape(shape, prs.slide_height)]
        page_number_shapes = []
        for shape in slide.shapes:
            if has_text(shape) and re.fullmatch(r"\d{1,2}", shape.text.strip()) and shape.top > prs.slide_height * 0.75:
                page_number_shapes.append({"text": shape.text.strip(), "bounds_in": shape_bounds(shape)})
        page_numbers.extend({"slide_index": idx, **item} for item in page_number_shapes)
        layouts.append(
            {
                "slide_index": idx,
                "layout_type": _template_layout_type(slide, prs, idx),
                "title": title,
                "text_box_count": sum(1 for shape in slide.shapes if has_text(shape)),
                "image_count": sum(1 for shape in slide.shapes if is_picture(shape)),
                "table_count": sum(1 for shape in slide.shapes if is_table(shape)),
                "title_bounds_in": shape_bounds(title_shapes[0]) if title_shapes else None,
            }
        )
    return {
        "ok": True,
        "template_pptx_path": str(path),
        "page_size": {
            "width_in": emu_to_inches(prs.slide_width),
            "height_in": emu_to_inches(prs.slide_height),
            "width_emu": int(prs.slide_width),
            "height_emu": int(prs.slide_height),
        },
        "common_background_or_fill_colors": _most_common(fills, 8),
        "title_bar_candidates": header_candidates[:10],
        "title_bar": header_candidates[0] if header_candidates else None,
        "page_title_style": {
            "fonts": _most_common(title_fonts, 5),
            "sizes_pt": _most_common(title_sizes, 5),
            "colors": _most_common(title_colors, 5),
            "sample_bounds": [item["title_bounds_in"] for item in layouts if item.get("title_bounds_in")][:5],
        },
        "body_style": {
            "fonts": _most_common(body_fonts, 5),
            "sizes_pt": _most_common(body_sizes, 5),
            "colors": _most_common(body_colors, 5),
        },
        "module_title_style": "Derived from bold non-title text runs in the template.",
        "layout_types": layouts,
        "logo_positions": logo_positions[:10],
        "image_classifications": image_classifications,
        "page_number_positions": page_numbers[:10],
        "emphasis_colors": _most_common(emphasis_colors, 5),
        "element_position_parameters": layouts,
    }


def _first_header_style(style: dict[str, Any]) -> dict[str, Any] | None:
    candidates = style.get("title_bar_candidates") or []
    return candidates[0] if candidates else None


def _first_title_bounds(style: dict[str, Any]) -> dict[str, float] | None:
    bounds = style.get("page_title_style", {}).get("sample_bounds") or []
    return bounds[0] if bounds else None


def apply_template_style(source_pptx_path: str, template_pptx_path: str, output_path: str) -> dict[str, Any]:
    """Conservatively apply template title-bar, font, margin, and logo style to a deck."""
    source = validate_pptx_path(source_pptx_path)
    template = validate_pptx_path(template_pptx_path)
    out_path = validate_output_path(output_path, source)
    style = inspect_template_style(str(template))
    prs = Presentation(str(source))
    header = _first_header_style(style)
    title_bounds = _first_title_bounds(style)
    title_color = (style.get("page_title_style", {}).get("colors") or [{"value": (255, 255, 255)}])[0]["value"]
    header_color = header["color"] if header else (29, 67, 113)
    logo_bounds = None
    if style.get("logo_positions"):
        logo_bounds = style["logo_positions"][0]["bounds_in"]
    modified = []
    for idx, slide in enumerate(prs.slides, start=1):
        if header and idx != 1:
            b = header["bounds_in"]
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches_to_emu(b["left"]), inches_to_emu(b["top"]), inches_to_emu(b["width"]), inches_to_emu(b["height"]))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*header_color)
            bar.line.color.rgb = RGBColor(*header_color)
            _send_to_back(bar)
        title_shape = None
        for shape in slide.shapes:
            if _is_title_shape(shape, prs.slide_height):
                title_shape = shape
                break
        if title_shape:
            if title_bounds and idx != 1:
                title_shape.left = inches_to_emu(title_bounds["left"])
                title_shape.top = inches_to_emu(title_bounds["top"])
                title_shape.width = inches_to_emu(min(title_bounds["width"], emu_to_inches(prs.slide_width) - title_bounds["left"] - 2.0))
                title_shape.height = inches_to_emu(title_bounds["height"])
            for _paragraph, run in iter_text_runs(title_shape) or []:
                _ensure_rfonts(run, "Times New Roman", "黑体")
                run.font.bold = True
                if idx != 1:
                    run.font.color.rgb = RGBColor(*title_color) if title_color else RGBColor(255, 255, 255)
        for shape in slide.shapes:
            if has_text(shape) and shape is not title_shape:
                for _paragraph, run in iter_text_runs(shape) or []:
                    _ensure_rfonts(run, "Times New Roman", "宋体")
        if logo_bounds and idx != 1:
            try:
                logo_path = _logo_asset_path(None)
                for old in _right_top_pictures(slide, prs):
                    _delete_shape(old)
                slide.shapes.add_picture(str(logo_path), inches_to_emu(logo_bounds["left"]), inches_to_emu(logo_bounds["top"]), width=inches_to_emu(logo_bounds["width"]))
            except PptToolError:
                pass
        modified.append(idx)
    prs.save(str(out_path))
    return {
        "ok": True,
        "source_pptx_path": str(source),
        "template_pptx_path": str(template),
        "output_path": str(out_path),
        "modified_slide_indices": modified,
        "applied_header": header is not None,
        "applied_logo_asset": logo_bounds is not None and DEFAULT_LOGO_PATH.exists(),
    }


def _interesting_shapes(slide: Any) -> list[dict[str, Any]]:
    items = []
    for shape in slide.shapes:
        kind = None
        text = ""
        if has_text(shape):
            kind = "text"
            text = shape.text.strip()[:40]
        elif is_picture(shape):
            kind = "image"
        elif is_table(shape):
            kind = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and _shape_fill_rgb(shape):
            # Include visible non-text shapes because they can cover titles/images,
            # but skip common full-width title bars and section backgrounds.
            bounds = shape_bounds(shape)
            if not (bounds["top"] <= 0.25 and bounds["width"] >= 9 and bounds["height"] <= 1.4):
                kind = "shape"
        if kind:
            items.append({"shape": shape, "kind": kind, "text": text, "bounds": shape_bounds(shape)})
    return items


def _overlap_area(a: dict[str, float], b: dict[str, float]) -> float:
    ax2, ay2 = a["left"] + a["width"], a["top"] + a["height"]
    bx2, by2 = b["left"] + b["width"], b["top"] + b["height"]
    w = max(0.0, min(ax2, bx2) - max(a["left"], b["left"]))
    h = max(0.0, min(ay2, by2) - max(a["top"], b["top"]))
    return w * h


def _unique_md_path(pptx_path: Path, suffix: str) -> Path:
    base = pptx_path.with_name(pptx_path.stem + suffix)
    if not base.exists():
        return base
    for i in range(1, 100):
        candidate = pptx_path.with_name(f"{pptx_path.stem}{suffix[:-3]}_{i}.md")
        if not candidate.exists():
            return candidate
    raise PptToolError("Could not create a unique markdown report path.")


def check_layout_overlap(pptx_path: str, report_path: str | None = None) -> dict[str, Any]:
    """Check for overlap, out-of-bounds objects, title obstruction, and logo obstruction."""
    path = validate_pptx_path(pptx_path)
    prs = Presentation(str(path))
    issues = []
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)
    for idx, slide in enumerate(prs.slides, start=1):
        items = _interesting_shapes(slide)
        title_items = [item for item in items if item["kind"] == "text" and _is_title_shape(item["shape"], prs.slide_height)]
        logo_items = [{"shape": sh, "kind": "logo", "bounds": shape_bounds(sh), "text": ""} for sh in _right_top_pictures(slide, prs)]
        for shape in slide.shapes:
            if not is_picture(shape):
                continue
            picture_info = _classify_picture_shape(shape, prs)
            b = picture_info["bounds_in"]
            right_ratio = (b["left"] + b["width"]) / max(slide_w, 0.001)
            top_ratio = b["top"] / max(slide_h, 0.001)
            if picture_info["classification"] != "header_logo" and right_ratio >= 0.74 and top_ratio <= 0.25 and picture_info["area_ratio"] > 0.03:
                issues.append(
                    {
                        "slide_index": idx,
                        "type": "image/logo_candidate",
                        "problem": "suspected_large_image_misidentified_as_logo",
                        "bounds_in": b,
                        "suggestion": "Delete the mistaken logo-like large picture, use the fixed CUMT logo asset, and reapply template style without copying content images.",
                    }
                )
            if picture_info["classification"] == "header_logo" and picture_info["area_ratio"] > 0.03:
                issues.append(
                    {
                        "slide_index": idx,
                        "type": "logo",
                        "problem": "logo_size_exceeds_reasonable_limit",
                        "bounds_in": b,
                        "suggestion": "Normalize logo size from a trusted reference or fixed CUMT logo asset.",
                    }
                )
            for title in title_items:
                if picture_info["area_ratio"] > 0.05 and _overlap_area(b, title["bounds"]) > 0:
                    issues.append(
                        {
                            "slide_index": idx,
                            "type": "image/title",
                            "problem": "large_image_overlaps_title_area",
                            "bounds_in": {"image": b, "title": title["bounds"]},
                            "suggestion": "Move or shrink the picture; do not place content images in the title/header area.",
                        }
                    )
        for item in items + logo_items:
            b = item["bounds"]
            if b["left"] < -0.01 or b["top"] < -0.01 or b["left"] + b["width"] > slide_w + 0.01 or b["top"] + b["height"] > slide_h + 0.01:
                issues.append({"slide_index": idx, "type": item["kind"], "problem": "object_out_of_bounds", "bounds_in": b, "suggestion": "Move or scale the object inside slide bounds."})
        for i in range(len(items)):
            for j in range(i + 1, len(items)):
                a, b = items[i], items[j]
                if a["kind"] == "text" and b["kind"] == "text":
                    continue
                area = _overlap_area(a["bounds"], b["bounds"])
                min_area = min(a["bounds"]["width"] * a["bounds"]["height"], b["bounds"]["width"] * b["bounds"]["height"])
                if min_area > 0 and area / min_area > 0.08:
                    issues.append({"slide_index": idx, "type": f"{a['kind']}/{b['kind']}", "problem": "objects_overlap", "bounds_in": {"a": a["bounds"], "b": b["bounds"]}, "suggestion": "Reduce image size or move it away from text/table areas."})
        for logo in logo_items:
            for title in title_items:
                if _overlap_area(logo["bounds"], title["bounds"]) > 0:
                    issues.append({"slide_index": idx, "type": "logo/title", "problem": "logo_overlaps_title", "bounds_in": {"logo": logo["bounds"], "title": title["bounds"]}, "suggestion": "Move logo further right/up or shorten title width."})
    out = _path(report_path) if report_path else _unique_md_path(path, "_layout_report.md")
    if out.suffix.lower() != ".md":
        raise PptToolError("Layout report path must end with .md.")
    out.parent.mkdir(parents=True, exist_ok=True)
    lines = ["# Layout Overlap Report", "", f"- PPTX: `{path}`", f"- Issue count: {len(issues)}", "", "| Slide | Type | Problem | Suggestion |", "|---:|---|---|---|"]
    for issue in issues:
        lines.append(f"| {issue['slide_index']} | {issue['type']} | {issue['problem']} | {issue['suggestion']} |")
    out.write_text("\n".join(lines), encoding="utf-8-sig")
    return {"ok": True, "pptx_path": str(path), "issue_count": len(issues), "issues": issues, "report_path": str(out)}


def auto_fix_layout(pptx_path: str, output_path: str, report_path: str | None = None) -> dict[str, Any]:
    """Conservatively fix obvious image overlap/out-of-bounds issues without changing text."""
    input_path = validate_pptx_path(pptx_path)
    out_path = validate_output_path(output_path, input_path)
    prs = Presentation(str(input_path))
    fixes = []
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)
    for idx, slide in enumerate(prs.slides, start=1):
        text_items = [item for item in _interesting_shapes(slide) if item["kind"] in {"text", "table"}]
        for pic in [shape for shape in slide.shapes if is_picture(shape)]:
            old = shape_bounds(pic)
            pic_bounds = shape_bounds(pic)
            overlaps = any(_overlap_area(pic_bounds, item["bounds"]) / max(pic_bounds["width"] * pic_bounds["height"], 0.001) > 0.06 for item in text_items)
            out_of_bounds = pic_bounds["left"] < 0 or pic_bounds["top"] < 0 or pic_bounds["left"] + pic_bounds["width"] > slide_w or pic_bounds["top"] + pic_bounds["height"] > slide_h
            if overlaps or out_of_bounds:
                max_w = slide_w * 0.42
                max_h = slide_h * 0.55
                _scale_shape_keep_ratio(pic, int(max_w * EMU_PER_INCH), int(max_h * EMU_PER_INCH))
                pic.left = inches_to_emu(slide_w - emu_to_inches(pic.width) - 0.75)
                pic.top = inches_to_emu(max(0.95, (slide_h - emu_to_inches(pic.height)) / 2))
                fixes.append({"slide_index": idx, "name": pic.name, "old_bounds_in": old, "new_bounds_in": shape_bounds(pic)})
    prs.save(str(out_path))
    report = _path(report_path) if report_path else out_path.with_name(out_path.stem + "_layout_fix_report.md")
    report.parent.mkdir(parents=True, exist_ok=True)
    lines = ["# Auto Layout Fix Report", "", f"- Input: `{input_path}`", f"- Output: `{out_path}`", f"- Fix count: {len(fixes)}", "", "| Slide | Shape | Change |", "|---:|---|---|"]
    for fix in fixes:
        lines.append(f"| {fix['slide_index']} | {fix['name']} | {fix['old_bounds_in']} -> {fix['new_bounds_in']} |")
    report.write_text("\n".join(lines), encoding="utf-8-sig")
    return {"ok": True, "input_path": str(input_path), "output_path": str(out_path), "fix_count": len(fixes), "fixes": fixes, "report_path": str(report)}


def _background_rgb(background_color: str) -> tuple[int, int, int]:
    if background_color.lower() == "white":
        return (255, 255, 255)
    if background_color.lower() == "black":
        return (0, 0, 0)
    m = re.fullmatch(r"#?([0-9a-fA-F]{6})", background_color.strip())
    if not m:
        raise PptToolError("background_color must be 'white', 'black', or a hex color like #ffffff.")
    hex_value = m.group(1)
    return tuple(int(hex_value[i : i + 2], 16) for i in (0, 2, 4))


def _save_image_no_black(img: Image.Image, path: Path, bg_rgb: tuple[int, int, int]) -> dict[str, Any]:
    rgba = img.convert("RGBA")
    bg = Image.new("RGBA", rgba.size, (*bg_rgb, 255))
    bg.alpha_composite(rgba)
    rgb = bg.convert("RGB")
    rgb.save(path)
    pixels = list(rgb.getdata())
    dark_ratio = sum(1 for r, g, b in pixels if r < 35 and g < 35 and b < 35) / max(len(pixels), 1)
    return {"path": str(path), "width": rgb.width, "height": rgb.height, "dark_pixel_ratio": round(dark_ratio, 4)}


def extract_pdf_images_safe(pdf_path: str, output_dir: str, background_color: str = "white") -> dict[str, Any]:
    """Extract embedded PDF images and page renders using a white background to avoid black alpha fills."""
    pdf = validate_pdf_path(pdf_path)
    out_dir = _path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    image_dir = out_dir / "images"
    render_dir = out_dir / "page_renders"
    image_dir.mkdir(parents=True, exist_ok=True)
    render_dir.mkdir(parents=True, exist_ok=True)
    bg_rgb = _background_rgb(background_color)
    try:
        import fitz
    except Exception as exc:
        raise PptToolError(f"PyMuPDF/fitz is required for PDF extraction: {exc}") from exc
    doc = fitz.open(str(pdf))
    extracted = []
    rendered = []
    for page_no, page in enumerate(doc, start=1):
        for img_no, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.n >= 5:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            data = pix.tobytes("png")
            pil = Image.open(BytesIO(data))
            path = image_dir / f"p{page_no:03d}_img{img_no:02d}.png"
            info = _save_image_no_black(pil, path, bg_rgb)
            info.update({"page": page_no, "image_index": img_no, "method": "pymupdf_get_images"})
            extracted.append(info)
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        pil = Image.open(BytesIO(pix.tobytes("png")))
        render_path = render_dir / f"page_{page_no:03d}.png"
        info = _save_image_no_black(pil, render_path, bg_rgb)
        info.update({"page": page_no, "method": "pymupdf_render_page_white_background"})
        rendered.append(info)
    pdf2image_available = shutil.which("pdftoppm") is not None
    report_path = out_dir / "extract_pdf_images_safe_report.md"
    lines = [
        "# Safe PDF Image Extraction Report",
        "",
        f"- PDF: `{pdf}`",
        f"- Background color: `{background_color}`",
        f"- Extracted embedded images: {len(extracted)}",
        f"- Rendered pages: {len(rendered)}",
        f"- pdf2image/poppler available: {pdf2image_available}",
        "",
        "## Black Background Mitigation",
        "- Embedded images are converted to RGBA and composited onto the requested background color.",
        "- Page renders use PyMuPDF with `alpha=False` and are saved on a white background by default.",
        "- Dark-pixel ratios are reported for follow-up inspection.",
    ]
    report_path.write_text("\n".join(lines), encoding="utf-8-sig")
    inventory_path = out_dir / "image_inventory.json"
    inventory = {"pdf_path": str(pdf), "background_color": background_color, "extracted_images": extracted, "rendered_pages": rendered, "report_path": str(report_path)}
    inventory_path.write_text(json_dumps(inventory), encoding="utf-8")
    return {"ok": True, "pdf_path": str(pdf), "output_dir": str(out_dir), "image_count": len(extracted), "rendered_page_count": len(rendered), "inventory_path": str(inventory_path), "report_path": str(report_path), "extracted_images": extracted[:20], "rendered_pages": rendered[:10]}


def json_dumps(data: Any) -> str:
    import json

    return json.dumps(data, ensure_ascii=False, indent=2)
